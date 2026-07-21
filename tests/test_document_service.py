"""Tests for isolated Word and PowerPoint document generation."""

from __future__ import annotations

import base64
import zipfile
from io import BytesIO
from pathlib import Path

import pytest
from pydantic import ValidationError

import microsoft_agent.document_service as documents

CAPABILITIES = documents.get_document_capabilities()


def test_capabilities_report_optional_backends() -> None:
    """Capabilities are available without requiring optional Office packages."""
    capabilities = documents.get_document_capabilities()

    assert isinstance(capabilities.word_available, bool)
    assert isinstance(capabilities.powerpoint_available, bool)
    assert "python-docx" in capabilities.word_install_hint
    assert "python-pptx" in capabilities.powerpoint_install_hint


def test_request_models_reject_unsafe_names_and_jagged_tables() -> None:
    """Model validation blocks paths in filenames and malformed tables."""
    with pytest.raises(ValidationError):
        documents.WordDocumentRequest(filename="../report.docx")

    with pytest.raises(ValidationError):
        documents.PowerPointPresentationRequest(filename="slides.pptm")

    with pytest.raises(ValidationError):
        documents.WordTable(headers=["A", "B"], rows=[["one"]])

    with pytest.raises(ValidationError):
        documents.WordDocumentRequest(replacements={"not valid": "value"})


def test_artifact_options_require_unambiguous_file_delivery() -> None:
    """Output paths cannot be silently ignored or omitted for file delivery."""
    with pytest.raises(ValidationError):
        documents.ArtifactOptions(delivery="file")

    with pytest.raises(ValidationError):
        documents.ArtifactOptions(delivery="bytes", output_path="unused.docx")


@pytest.mark.asyncio
async def test_missing_word_dependency_has_actionable_error(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """A missing Word backend fails at call time, not module import time."""
    monkeypatch.setattr(documents, "_DocxDocument", None)
    service = documents.DocumentService(tmp_path)

    with pytest.raises(documents.OptionalDocumentDependencyError, match="python-docx"):
        await service.generate_word_document(documents.WordDocumentRequest())


@pytest.mark.asyncio
async def test_missing_powerpoint_dependency_has_actionable_error(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """A missing PowerPoint backend gives an actionable optional-install hint."""
    monkeypatch.setattr(documents, "_PowerPointPresentation", None)
    service = documents.DocumentService(tmp_path)

    with pytest.raises(documents.OptionalDocumentDependencyError, match="python-pptx"):
        await service.generate_powerpoint_presentation(
            documents.PowerPointPresentationRequest()
        )


@pytest.mark.skipif(
    not CAPABILITIES.word_available, reason="python-docx is not installed"
)
@pytest.mark.asyncio
async def test_generate_word_bytes_for_graph_upload(tmp_path: Path) -> None:
    """Word generation returns valid bytes, metadata, and a stable digest."""
    from docx import Document

    service = documents.DocumentService(tmp_path / "artifacts")
    request = documents.WordDocumentRequest(
        filename="status.docx",
        title="Project Status",
        paragraphs=[
            documents.WordParagraph(
                runs=[
                    documents.WordTextRun(text="On track", bold=True),
                    documents.WordTextRun(text=" for launch."),
                ]
            )
        ],
        tables=[
            documents.WordTable(headers=["Owner", "State"], rows=[["Alex", "Ready"]])
        ],
        metadata=documents.DocumentMetadata(author="Microsoft Agent"),
        artifact=documents.ArtifactOptions(delivery="bytes"),
    )

    artifact = await service.generate_word_document(request)

    assert artifact.content is not None
    assert artifact.content.startswith(b"PK")
    assert artifact.content_base64 is None
    assert artifact.upload_bytes() == artifact.content
    parsed = Document(BytesIO(artifact.content))
    assert parsed.paragraphs[0].text == "Project Status"
    assert parsed.paragraphs[1].text == "On track for launch."
    assert parsed.tables[0].cell(1, 1).text == "Ready"
    assert parsed.core_properties.author == "Microsoft Agent"


@pytest.mark.skipif(
    not CAPABILITIES.word_available, reason="python-docx is not installed"
)
@pytest.mark.asyncio
async def test_word_template_replacement_preserves_split_run_formatting(
    tmp_path: Path,
) -> None:
    """Template tokens spanning runs are replaced without flattening all runs."""
    from docx import Document

    template_root = tmp_path / "templates"
    template_root.mkdir()
    template_path = template_root / "letter.docx"
    template = Document()
    paragraph = template.add_paragraph()
    first = paragraph.add_run("Hello {{na")
    first.bold = True
    paragraph.add_run("me}}!")
    template.save(template_path)

    service = documents.DocumentService(
        tmp_path / "artifacts", template_root=template_root
    )
    artifact = await service.generate_word_document(
        documents.WordDocumentRequest(
            filename="letter.docx",
            template_path="letter.docx",
            replacements={"name": "Sample User"},
            artifact=documents.ArtifactOptions(delivery="bytes"),
        )
    )

    assert artifact.content is not None
    generated = Document(BytesIO(artifact.content))
    generated_paragraph = generated.paragraphs[0]
    assert generated_paragraph.text == "Hello Sample User!"
    assert generated_paragraph.runs[0].bold is True
    assert generated_paragraph.runs[-1].text == "!"


@pytest.mark.skipif(
    not CAPABILITIES.word_available, reason="python-docx is not installed"
)
@pytest.mark.asyncio
async def test_word_file_delivery_is_contained_and_exclusive(tmp_path: Path) -> None:
    """Artifacts write beneath their root and do not overwrite by default."""
    artifact_root = tmp_path / "artifacts"
    service = documents.DocumentService(artifact_root)
    request = documents.WordDocumentRequest(
        title="Saved",
        artifact=documents.ArtifactOptions(
            delivery="file", output_path="reports/status.docx"
        ),
    )

    artifact = await service.generate_word_document(request)

    expected = (artifact_root / "reports/status.docx").resolve()
    assert artifact.path == expected
    assert expected.read_bytes().startswith(b"PK")
    with pytest.raises(documents.DocumentServiceError):
        artifact.upload_bytes()
    with pytest.raises(documents.ArtifactExistsError):
        await service.generate_word_document(request)


@pytest.mark.skipif(
    not CAPABILITIES.word_available, reason="python-docx is not installed"
)
@pytest.mark.asyncio
async def test_file_delivery_rejects_parent_traversal(tmp_path: Path) -> None:
    """Relative traversal cannot write an artifact outside its configured root."""
    service = documents.DocumentService(tmp_path / "artifacts")
    request = documents.WordDocumentRequest(
        artifact=documents.ArtifactOptions(
            delivery="file", output_path="../escaped.docx"
        )
    )

    with pytest.raises(documents.DocumentPathError, match="configured root"):
        await service.generate_word_document(request)
    assert not (tmp_path / "escaped.docx").exists()


@pytest.mark.skipif(
    not CAPABILITIES.word_available, reason="python-docx is not installed"
)
@pytest.mark.asyncio
async def test_file_delivery_rejects_symlink_escape(tmp_path: Path) -> None:
    """Symlink resolution cannot redirect artifact writes outside the root."""
    artifact_root = tmp_path / "artifacts"
    outside = tmp_path / "outside"
    artifact_root.mkdir()
    outside.mkdir()
    try:
        (artifact_root / "link").symlink_to(outside, target_is_directory=True)
    except OSError:
        pytest.skip("directory symlinks are unavailable on this platform")
    service = documents.DocumentService(artifact_root)
    request = documents.WordDocumentRequest(
        artifact=documents.ArtifactOptions(
            delivery="file", output_path="link/escaped.docx"
        )
    )

    with pytest.raises(documents.DocumentPathError, match="configured root"):
        await service.generate_word_document(request)
    assert not (outside / "escaped.docx").exists()


@pytest.mark.skipif(
    not CAPABILITIES.word_available, reason="python-docx is not installed"
)
@pytest.mark.asyncio
async def test_generated_artifact_size_limit_is_enforced(tmp_path: Path) -> None:
    """Large generated packages are rejected before delivery."""
    service = documents.DocumentService(tmp_path, max_artifact_bytes=10)

    with pytest.raises(documents.ArtifactTooLargeError):
        await service.generate_word_document(documents.WordDocumentRequest())


def test_template_validation_blocks_embedded_active_content(tmp_path: Path) -> None:
    """Office packages containing embedded objects are denied by default."""
    template = tmp_path / "unsafe.docx"
    with zipfile.ZipFile(template, "w") as package:
        package.writestr("[Content_Types].xml", "<Types />")
        package.writestr("word/embeddings/object1.bin", b"payload")
    service = documents.DocumentService(tmp_path, template_root=tmp_path)

    with pytest.raises(documents.DocumentTemplateError, match="embedded content"):
        service._validate_office_package(template)


def test_template_validation_blocks_external_relationships(tmp_path: Path) -> None:
    """Templates cannot retain network-loading relationships by default."""
    template = tmp_path / "external.pptx"
    relationships = b"""<?xml version="1.0"?>
    <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
      <Relationship Id="rId1" Type="example" Target="https://example.com/a"
                    TargetMode="External" />
    </Relationships>"""
    with zipfile.ZipFile(template, "w") as package:
        package.writestr("[Content_Types].xml", "<Types />")
        package.writestr("ppt/_rels/presentation.xml.rels", relationships)
    service = documents.DocumentService(tmp_path, template_root=tmp_path)

    with pytest.raises(documents.DocumentTemplateError, match="external relationship"):
        service._validate_office_package(template)


@pytest.mark.skipif(
    not CAPABILITIES.powerpoint_available, reason="python-pptx is not installed"
)
@pytest.mark.asyncio
async def test_generate_powerpoint_base64_for_json_transport(tmp_path: Path) -> None:
    """PowerPoint generation returns valid base64 with expected slide content."""
    from pptx import Presentation

    service = documents.DocumentService(tmp_path / "artifacts")
    request = documents.PowerPointPresentationRequest(
        filename="briefing.pptx",
        title_slide=documents.PowerPointTitleSlide(
            title="Quarterly Briefing", subtitle="Microsoft Agent"
        ),
        slides=[
            documents.PowerPointSlide(
                title="Highlights",
                paragraphs=[
                    documents.PowerPointParagraph(text="Revenue increased", bold=True),
                    documents.PowerPointParagraph(text="Launch ready", level=1),
                ],
            )
        ],
        artifact=documents.ArtifactOptions(delivery="base64"),
    )

    artifact = await service.generate_powerpoint_presentation(request)

    assert artifact.content is None
    assert artifact.content_base64 is not None
    decoded = base64.b64decode(artifact.content_base64, validate=True)
    assert decoded.startswith(b"PK")
    assert artifact.upload_bytes() == decoded
    presentation = Presentation(BytesIO(decoded))
    assert len(presentation.slides) == 2
    first_text = [
        shape.text for shape in presentation.slides[0].shapes if shape.has_text_frame
    ]
    second_text = [
        shape.text for shape in presentation.slides[1].shapes if shape.has_text_frame
    ]
    assert any("Quarterly Briefing" in text for text in first_text)
    assert any("Highlights" in text for text in second_text)
    assert any("Revenue increased" in text for text in second_text)


@pytest.mark.skipif(
    not CAPABILITIES.powerpoint_available, reason="python-pptx is not installed"
)
@pytest.mark.asyncio
async def test_powerpoint_rejects_missing_layout(tmp_path: Path) -> None:
    """Requests receive a clear error when a template lacks a selected layout."""
    service = documents.DocumentService(tmp_path)
    request = documents.PowerPointPresentationRequest(
        slides=[documents.PowerPointSlide(title="No layout", layout_index=99)]
    )

    with pytest.raises(documents.DocumentGenerationError, match="layout index 99"):
        await service.generate_powerpoint_presentation(request)
