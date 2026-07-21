"""Safe, asynchronous Word and PowerPoint document generation.

The service deliberately keeps Microsoft Graph concerns out of document creation.
Callers can request raw bytes for an SDK upload, base64 for JSON transport, or a
file confined to a configured artifact directory.  ``python-docx`` and
``python-pptx`` are optional: importing this module never requires either one.
"""

from __future__ import annotations

import asyncio
import base64
import hashlib
import os
import re
import tempfile
import zipfile
from bisect import bisect_right
from enum import StrEnum
from io import BytesIO
from pathlib import Path, PurePosixPath
from typing import Any, Self
from xml.etree import ElementTree

from pydantic import BaseModel, ConfigDict, Field, field_validator, model_validator

try:
    from docx import Document as _DocxDocument
    from docx.enum.text import WD_ALIGN_PARAGRAPH as _WordAlignment
    from docx.shared import Pt as _DocxPoints
except ImportError:  # pragma: no cover - exercised by monkeypatch in tests
    _DocxDocument = None  # type: ignore[assignment]
    _WordAlignment = None  # type: ignore[assignment,misc]
    _DocxPoints = None  # type: ignore[assignment,misc]

try:
    from pptx import Presentation as _PowerPointPresentation
    from pptx.util import Pt as _PowerPointPoints
except ImportError:  # pragma: no cover - exercised by monkeypatch in tests
    _PowerPointPresentation = None  # type: ignore[assignment]
    _PowerPointPoints = None  # type: ignore[assignment,misc]


_PLACEHOLDER_NAME = re.compile(r"^[A-Za-z][A-Za-z0-9_.-]{0,127}$")
_MIME_TYPES = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class DocumentServiceError(RuntimeError):
    """Base error for document validation, generation, and delivery failures."""


class OptionalDocumentDependencyError(DocumentServiceError):
    """Raised when the library needed for a document type is not installed."""


class DocumentPathError(DocumentServiceError):
    """Raised when a template or artifact escapes its configured directory."""


class DocumentTemplateError(DocumentServiceError):
    """Raised when an Office template is missing, unsafe, or malformed."""


class DocumentGenerationError(DocumentServiceError):
    """Raised when an Office library cannot generate the requested document."""


class ArtifactExistsError(DocumentServiceError):
    """Raised when an artifact exists and the request did not allow overwrite."""


class ArtifactTooLargeError(DocumentServiceError):
    """Raised when generated content exceeds the configured output limit."""


class ArtifactDelivery(StrEnum):
    """Supported ways to return a generated Office document."""

    BYTES = "bytes"
    BASE64 = "base64"
    FILE = "file"


class ParagraphAlignment(StrEnum):
    """Portable paragraph alignment options for generated Word content."""

    LEFT = "left"
    CENTER = "center"
    RIGHT = "right"
    JUSTIFY = "justify"


class DocumentMetadata(BaseModel):
    """Optional core properties written into the generated Office package."""

    model_config = ConfigDict(extra="forbid")

    title: str | None = Field(default=None, max_length=255)
    subject: str | None = Field(default=None, max_length=255)
    author: str | None = Field(default=None, max_length=255)
    category: str | None = Field(default=None, max_length=255)
    keywords: str | None = Field(default=None, max_length=255)
    comments: str | None = Field(default=None, max_length=2_000)


class ArtifactOptions(BaseModel):
    """Select document delivery and control filesystem overwrite behavior."""

    model_config = ConfigDict(extra="forbid")

    delivery: ArtifactDelivery = ArtifactDelivery.BASE64
    output_path: Path | None = None
    overwrite: bool = False

    @model_validator(mode="after")
    def validate_delivery_path(self) -> Self:
        """Require a path only for file delivery to prevent ambiguous output."""
        if self.delivery is ArtifactDelivery.FILE and self.output_path is None:
            raise ValueError("output_path is required for file delivery")
        if self.delivery is not ArtifactDelivery.FILE and self.output_path is not None:
            raise ValueError("output_path is only valid for file delivery")
        return self


class GeneratedArtifact(BaseModel):
    """A generated document ready for Graph upload, JSON, or local use."""

    model_config = ConfigDict(extra="forbid", ser_json_bytes="base64")

    document_type: str
    filename: str
    content_type: str
    size_bytes: int = Field(ge=0)
    sha256: str = Field(min_length=64, max_length=64)
    delivery: ArtifactDelivery
    path: Path | None = None
    content: bytes | None = Field(default=None, repr=False)
    content_base64: str | None = Field(default=None, repr=False)

    def upload_bytes(self) -> bytes:
        """Return in-memory content suitable for a Microsoft Graph upload."""
        if self.content is not None:
            return self.content
        if self.content_base64 is not None:
            return base64.b64decode(self.content_base64, validate=True)
        raise DocumentServiceError(
            "file-delivered artifacts must be read explicitly by their owner"
        )


class DocumentCapabilities(BaseModel):
    """Report which optional Office generation backends are available."""

    word_available: bool
    powerpoint_available: bool
    word_install_hint: str = "pip install python-docx"
    powerpoint_install_hint: str = "pip install python-pptx"


class WordTextRun(BaseModel):
    """A formatted span inside a generated Word paragraph."""

    model_config = ConfigDict(extra="forbid")

    text: str = Field(max_length=100_000)
    bold: bool = False
    italic: bool = False
    underline: bool = False
    font_size_points: float | None = Field(default=None, ge=1, le=400)


class WordParagraph(BaseModel):
    """A Word paragraph containing plain text and/or formatted runs."""

    model_config = ConfigDict(extra="forbid")

    text: str = Field(default="", max_length=100_000)
    runs: list[WordTextRun] = Field(default_factory=list, max_length=1_000)
    style: str | None = Field(default=None, max_length=255)
    alignment: ParagraphAlignment | None = None
    page_break_before: bool = False
    keep_with_next: bool = False


class WordTable(BaseModel):
    """A rectangular Word table with an optional formatted header row."""

    model_config = ConfigDict(extra="forbid")

    headers: list[str] = Field(default_factory=list, max_length=50)
    rows: list[list[str]] = Field(default_factory=list, max_length=1_000)
    style: str | None = Field(default=None, max_length=255)

    @model_validator(mode="after")
    def validate_rectangular_table(self) -> Self:
        """Reject empty, oversized, or jagged tables before generation."""
        if not self.headers and not self.rows:
            raise ValueError("a table requires headers or at least one row")
        width = len(self.headers) if self.headers else len(self.rows[0])
        if width == 0:
            raise ValueError("table rows must contain at least one cell")
        if width > 50:
            raise ValueError("tables may contain at most 50 columns")
        for row in self.rows:
            if len(row) != width:
                raise ValueError("all table rows must have the same width")
            if any(len(cell) > 100_000 for cell in row):
                raise ValueError("table cells may contain at most 100000 characters")
        if any(len(cell) > 100_000 for cell in self.headers):
            raise ValueError("table cells may contain at most 100000 characters")
        return self


class _DocumentRequest(BaseModel):
    """Shared validated fields for Word and PowerPoint requests."""

    model_config = ConfigDict(extra="forbid")

    template_path: Path | None = None
    replacements: dict[str, str] = Field(default_factory=dict, max_length=200)
    metadata: DocumentMetadata = Field(default_factory=DocumentMetadata)
    artifact: ArtifactOptions = Field(default_factory=ArtifactOptions)

    @field_validator("replacements")
    @classmethod
    def validate_replacements(cls, value: dict[str, str]) -> dict[str, str]:
        """Constrain ``{{name}}`` placeholder names and replacement sizes."""
        for name, replacement in value.items():
            if not _PLACEHOLDER_NAME.fullmatch(name):
                raise ValueError(
                    "replacement names must start with a letter and contain only "
                    "letters, numbers, dots, underscores, or hyphens"
                )
            if len(replacement) > 1_000_000:
                raise ValueError(
                    "replacement values may contain at most 1000000 characters"
                )
        return value


class WordDocumentRequest(_DocumentRequest):
    """Validated inputs for creating or filling a ``.docx`` document."""

    filename: str = "document.docx"
    title: str | None = Field(default=None, max_length=10_000)
    paragraphs: list[WordParagraph] = Field(default_factory=list, max_length=5_000)
    tables: list[WordTable] = Field(default_factory=list, max_length=200)

    @field_validator("filename")
    @classmethod
    def validate_filename(cls, value: str) -> str:
        """Require a portable basename with the Word OOXML extension."""
        return _safe_filename(value, ".docx")

    @field_validator("template_path")
    @classmethod
    def validate_template_extension(cls, value: Path | None) -> Path | None:
        """Accept only non-macro Word OOXML templates."""
        return _template_extension(value, ".docx")


class PowerPointParagraph(BaseModel):
    """A paragraph or bullet placed in a PowerPoint content placeholder."""

    model_config = ConfigDict(extra="forbid")

    text: str = Field(max_length=100_000)
    level: int = Field(default=0, ge=0, le=8)
    bold: bool = False
    italic: bool = False
    font_size_points: float | None = Field(default=None, ge=1, le=400)


class PowerPointTitleSlide(BaseModel):
    """An optional opening title slide for a PowerPoint presentation."""

    model_config = ConfigDict(extra="forbid")

    title: str = Field(min_length=1, max_length=10_000)
    subtitle: str | None = Field(default=None, max_length=20_000)
    layout_index: int = Field(default=0, ge=0, le=99)


class PowerPointSlide(BaseModel):
    """A title-and-content slide appended to a PowerPoint presentation."""

    model_config = ConfigDict(extra="forbid")

    title: str = Field(default="", max_length=10_000)
    paragraphs: list[PowerPointParagraph] = Field(
        default_factory=list, max_length=1_000
    )
    layout_index: int = Field(default=1, ge=0, le=99)


class PowerPointPresentationRequest(_DocumentRequest):
    """Validated inputs for creating or filling a ``.pptx`` presentation."""

    filename: str = "presentation.pptx"
    title_slide: PowerPointTitleSlide | None = None
    slides: list[PowerPointSlide] = Field(default_factory=list, max_length=1_000)

    @field_validator("filename")
    @classmethod
    def validate_filename(cls, value: str) -> str:
        """Require a portable basename with the PowerPoint OOXML extension."""
        return _safe_filename(value, ".pptx")

    @field_validator("template_path")
    @classmethod
    def validate_template_extension(cls, value: Path | None) -> Path | None:
        """Accept only non-macro PowerPoint OOXML templates."""
        return _template_extension(value, ".pptx")


def get_document_capabilities() -> DocumentCapabilities:
    """Return optional dependency availability without importing server code."""
    return DocumentCapabilities(
        word_available=_DocxDocument is not None,
        powerpoint_available=_PowerPointPresentation is not None,
    )


class DocumentService:
    """Generate Office documents inside explicit template and artifact roots.

    ``template_root`` is a read trust boundary and ``artifact_root`` is a write
    boundary.  Absolute paths are accepted only when they remain within the
    corresponding root after symlink resolution.
    """

    def __init__(
        self,
        artifact_root: str | Path,
        template_root: str | Path | None = None,
        *,
        max_template_bytes: int = 50 * 1024 * 1024,
        max_uncompressed_template_bytes: int = 250 * 1024 * 1024,
        max_artifact_bytes: int = 100 * 1024 * 1024,
        max_template_entries: int = 10_000,
        allow_external_relationships: bool = False,
        allow_embedded_objects: bool = False,
    ) -> None:
        """Configure path roots, package limits, and active-content policy."""
        if (
            min(
                max_template_bytes,
                max_uncompressed_template_bytes,
                max_artifact_bytes,
                max_template_entries,
            )
            <= 0
        ):
            raise ValueError("document service size and entry limits must be positive")
        self.artifact_root = Path(artifact_root).expanduser().resolve(strict=False)
        template_base = template_root if template_root is not None else artifact_root
        self.template_root = Path(template_base).expanduser().resolve(strict=False)
        self.max_template_bytes = max_template_bytes
        self.max_uncompressed_template_bytes = max_uncompressed_template_bytes
        self.max_artifact_bytes = max_artifact_bytes
        self.max_template_entries = max_template_entries
        self.allow_external_relationships = allow_external_relationships
        self.allow_embedded_objects = allow_embedded_objects

    async def generate_word_document(
        self, request: WordDocumentRequest
    ) -> GeneratedArtifact:
        """Create or fill a Word document without blocking the event loop."""
        if _DocxDocument is None:
            raise OptionalDocumentDependencyError(
                "Word generation requires the optional 'python-docx' package"
            )
        return await asyncio.to_thread(self._generate_word_sync, request)

    async def generate_powerpoint_presentation(
        self, request: PowerPointPresentationRequest
    ) -> GeneratedArtifact:
        """Create or fill a PowerPoint file without blocking the event loop."""
        if _PowerPointPresentation is None:
            raise OptionalDocumentDependencyError(
                "PowerPoint generation requires the optional 'python-pptx' package"
            )
        return await asyncio.to_thread(self._generate_powerpoint_sync, request)

    def _generate_word_sync(self, request: WordDocumentRequest) -> GeneratedArtifact:
        try:
            template = self._prepare_template(request.template_path, ".docx")
            output_path = self._prepare_output(request.artifact, ".docx")
            document = _DocxDocument(str(template)) if template else _DocxDocument()
            replacement_plan = _replacement_plan(request.replacements)
            if replacement_plan is not None:
                _replace_word_document(document, replacement_plan)
            _set_core_properties(document.core_properties, request.metadata)
            _append_word_content(document, request)
            buffer = BytesIO()
            document.save(buffer)
            content = buffer.getvalue()
            return self._deliver(
                content,
                request.filename,
                ".docx",
                request.artifact,
                output_path,
            )
        except DocumentServiceError:
            raise
        except Exception as exc:
            raise DocumentGenerationError("Word document generation failed") from exc

    def _generate_powerpoint_sync(
        self, request: PowerPointPresentationRequest
    ) -> GeneratedArtifact:
        try:
            template = self._prepare_template(request.template_path, ".pptx")
            output_path = self._prepare_output(request.artifact, ".pptx")
            presentation = (
                _PowerPointPresentation(str(template))
                if template
                else _PowerPointPresentation()
            )
            replacement_plan = _replacement_plan(request.replacements)
            if replacement_plan is not None:
                _replace_powerpoint_presentation(presentation, replacement_plan)
            _set_core_properties(presentation.core_properties, request.metadata)
            _append_powerpoint_content(presentation, request)
            buffer = BytesIO()
            presentation.save(buffer)
            content = buffer.getvalue()
            return self._deliver(
                content,
                request.filename,
                ".pptx",
                request.artifact,
                output_path,
            )
        except DocumentServiceError:
            raise
        except Exception as exc:
            raise DocumentGenerationError(
                "PowerPoint presentation generation failed"
            ) from exc

    def _prepare_template(self, requested: Path | None, extension: str) -> Path | None:
        if requested is None:
            return None
        template = self._contained_path(
            self.template_root, requested, kind="template", must_exist=True
        )
        if not template.is_file():
            raise DocumentTemplateError("the requested template is not a regular file")
        if template.suffix.lower() != extension:
            raise DocumentTemplateError(f"templates must use the {extension} extension")
        self._validate_office_package(template)
        return template

    def _prepare_output(self, options: ArtifactOptions, extension: str) -> Path | None:
        if options.delivery is not ArtifactDelivery.FILE:
            return None
        assert options.output_path is not None
        output = self._contained_path(
            self.artifact_root,
            options.output_path,
            kind="artifact",
            must_exist=False,
        )
        if output.suffix.lower() != extension:
            raise DocumentPathError(
                f"artifact paths must use the {extension} extension"
            )
        output.parent.mkdir(parents=True, exist_ok=True)
        safe_parent = output.parent.resolve(strict=True)
        self._assert_contained(self.artifact_root, safe_parent, "artifact")
        output = safe_parent / output.name
        if output.exists() and not options.overwrite:
            raise ArtifactExistsError(f"artifact already exists: {output.name}")
        if output.exists() and not output.is_file():
            raise DocumentPathError("artifact path does not identify a regular file")
        return output

    def _contained_path(
        self,
        root: Path,
        requested: Path,
        *,
        kind: str,
        must_exist: bool,
    ) -> Path:
        raw = str(requested)
        if "\x00" in raw:
            raise DocumentPathError(f"{kind} path contains a null byte")
        candidate = requested if requested.is_absolute() else root / requested
        try:
            resolved = candidate.resolve(strict=must_exist)
        except FileNotFoundError as exc:
            if kind == "template":
                raise DocumentTemplateError(
                    "the requested template does not exist"
                ) from exc
            raise DocumentPathError(f"{kind} path does not exist") from exc
        self._assert_contained(root, resolved, kind)
        return resolved

    @staticmethod
    def _assert_contained(root: Path, candidate: Path, kind: str) -> None:
        canonical_root = root.resolve(strict=False)
        try:
            candidate.relative_to(canonical_root)
        except ValueError as exc:
            raise DocumentPathError(
                f"{kind} path must remain inside its configured root"
            ) from exc

    def _validate_office_package(self, template: Path) -> None:
        try:
            if template.stat().st_size > self.max_template_bytes:
                raise DocumentTemplateError(
                    "template exceeds the compressed size limit"
                )
            with zipfile.ZipFile(template) as package:
                entries = package.infolist()
                if len(entries) > self.max_template_entries:
                    raise DocumentTemplateError(
                        "template contains too many package parts"
                    )
                unpacked_size = 0
                seen: set[str] = set()
                for entry in entries:
                    name = entry.filename
                    member = PurePosixPath(name)
                    if (
                        member.is_absolute()
                        or ".." in member.parts
                        or "\\" in name
                        or name in seen
                    ):
                        raise DocumentTemplateError(
                            "template contains an unsafe package part path"
                        )
                    seen.add(name)
                    if entry.flag_bits & 0x1:
                        raise DocumentTemplateError(
                            "encrypted Office templates are not supported"
                        )
                    unpacked_size += entry.file_size
                    if unpacked_size > self.max_uncompressed_template_bytes:
                        raise DocumentTemplateError(
                            "template exceeds the uncompressed size limit"
                        )
                    lowered = f"/{name.lower()}"
                    if not self.allow_embedded_objects and (
                        lowered.endswith("/vbaproject.bin")
                        or "/activex/" in lowered
                        or "/embeddings/" in lowered
                    ):
                        raise DocumentTemplateError(
                            "template contains blocked active or embedded content"
                        )
                    if (
                        not self.allow_external_relationships
                        and lowered.endswith(".rels")
                        and _contains_external_relationship(package.read(entry))
                    ):
                        raise DocumentTemplateError(
                            "template contains a blocked external relationship"
                        )
        except zipfile.BadZipFile as exc:
            raise DocumentTemplateError(
                "template is not a valid Office Open XML package"
            ) from exc
        except OSError as exc:
            raise DocumentTemplateError("template could not be read") from exc

    def _deliver(
        self,
        content: bytes,
        filename: str,
        extension: str,
        options: ArtifactOptions,
        output_path: Path | None,
    ) -> GeneratedArtifact:
        size = len(content)
        if size > self.max_artifact_bytes:
            raise ArtifactTooLargeError("generated document exceeds the artifact limit")
        digest = hashlib.sha256(content).hexdigest()
        common: dict[str, Any] = {
            "document_type": "word" if extension == ".docx" else "powerpoint",
            "filename": output_path.name if output_path else filename,
            "content_type": _MIME_TYPES[extension],
            "size_bytes": size,
            "sha256": digest,
            "delivery": options.delivery,
        }
        if options.delivery is ArtifactDelivery.BYTES:
            return GeneratedArtifact(**common, content=content)
        if options.delivery is ArtifactDelivery.BASE64:
            encoded = base64.b64encode(content).decode("ascii")
            return GeneratedArtifact(**common, content_base64=encoded)
        assert output_path is not None
        _atomic_write(output_path, content, overwrite=options.overwrite)
        return GeneratedArtifact(**common, path=output_path)


async def generate_word_document(
    request: WordDocumentRequest,
    *,
    artifact_root: str | Path,
    template_root: str | Path | None = None,
) -> GeneratedArtifact:
    """Convenience API for one-off asynchronous Word generation."""
    service = DocumentService(artifact_root, template_root)
    return await service.generate_word_document(request)


async def generate_powerpoint_presentation(
    request: PowerPointPresentationRequest,
    *,
    artifact_root: str | Path,
    template_root: str | Path | None = None,
) -> GeneratedArtifact:
    """Convenience API for one-off asynchronous PowerPoint generation."""
    service = DocumentService(artifact_root, template_root)
    return await service.generate_powerpoint_presentation(request)


def _safe_filename(value: str, extension: str) -> str:
    if (
        not value
        or value in {".", ".."}
        or Path(value).name != value
        or "/" in value
        or "\\" in value
        or "\x00" in value
    ):
        raise ValueError("filename must be a non-empty portable basename")
    if not value.lower().endswith(extension):
        raise ValueError(f"filename must use the {extension} extension")
    return value


def _template_extension(value: Path | None, extension: str) -> Path | None:
    if value is not None and value.suffix.lower() != extension:
        raise ValueError(f"template_path must use the {extension} extension")
    return value


def _contains_external_relationship(content: bytes) -> bool:
    try:
        root = ElementTree.fromstring(content)
    except ElementTree.ParseError as exc:
        raise DocumentTemplateError(
            "template contains malformed package relationships"
        ) from exc
    return any(
        element.attrib.get("TargetMode", "").lower() == "external"
        for element in root.iter()
    )


def _replacement_plan(
    replacements: dict[str, str],
) -> tuple[re.Pattern[str], dict[str, str]] | None:
    if not replacements:
        return None
    token_values = {f"{{{{{name}}}}}": value for name, value in replacements.items()}
    alternatives = sorted(token_values, key=len, reverse=True)
    pattern = re.compile("|".join(re.escape(token) for token in alternatives))
    return pattern, token_values


def _replace_runs(runs: Any, plan: tuple[re.Pattern[str], dict[str, str]]) -> None:
    run_list = list(runs)
    if not run_list:
        return
    texts = [run.text or "" for run in run_list]
    full_text = "".join(texts)
    pattern, values = plan
    matches = list(pattern.finditer(full_text))
    if not matches:
        return
    starts: list[int] = []
    ends: list[int] = []
    cursor = 0
    for text in texts:
        starts.append(cursor)
        cursor += len(text)
        ends.append(cursor)
    output: list[list[str]] = [[] for _ in run_list]

    def copy_segment(start: int, end: int) -> None:
        position = start
        while position < end:
            index = bisect_right(ends, position)
            run_end = min(ends[index], end)
            output[index].append(full_text[position:run_end])
            position = run_end

    cursor = 0
    for match in matches:
        copy_segment(cursor, match.start())
        index = bisect_right(ends, match.start())
        if index >= len(run_list):
            index = len(run_list) - 1
        output[index].append(values[match.group(0)])
        cursor = match.end()
    copy_segment(cursor, len(full_text))
    for run, parts in zip(run_list, output, strict=True):
        run.text = "".join(parts)


def _replace_word_document(document: Any, plan: Any) -> None:
    for paragraph in document.paragraphs:
        _replace_runs(paragraph.runs, plan)
    for table in document.tables:
        _replace_word_table(table, plan)
    for section in document.sections:
        containers = (
            section.header,
            section.first_page_header,
            section.even_page_header,
            section.footer,
            section.first_page_footer,
            section.even_page_footer,
        )
        for container in containers:
            for paragraph in container.paragraphs:
                _replace_runs(paragraph.runs, plan)
            for table in container.tables:
                _replace_word_table(table, plan)


def _replace_word_table(table: Any, plan: Any) -> None:
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                _replace_runs(paragraph.runs, plan)
            for nested_table in cell.tables:
                _replace_word_table(nested_table, plan)


def _set_core_properties(properties: Any, metadata: DocumentMetadata) -> None:
    for name in ("title", "subject", "author", "category", "keywords", "comments"):
        value = getattr(metadata, name)
        if value is not None:
            setattr(properties, name, value)


def _append_word_content(document: Any, request: WordDocumentRequest) -> None:
    if request.title is not None:
        document.add_heading(request.title, level=0)
    alignments = {}
    if _WordAlignment is not None:
        alignments = {
            ParagraphAlignment.LEFT: _WordAlignment.LEFT,
            ParagraphAlignment.CENTER: _WordAlignment.CENTER,
            ParagraphAlignment.RIGHT: _WordAlignment.RIGHT,
            ParagraphAlignment.JUSTIFY: _WordAlignment.JUSTIFY,
        }
    for paragraph_spec in request.paragraphs:
        paragraph = document.add_paragraph(style=paragraph_spec.style)
        if paragraph_spec.text:
            paragraph.add_run(paragraph_spec.text)
        for span in paragraph_spec.runs:
            run = paragraph.add_run(span.text)
            run.bold = span.bold
            run.italic = span.italic
            run.underline = span.underline
            if span.font_size_points is not None and _DocxPoints is not None:
                run.font.size = _DocxPoints(span.font_size_points)
        if paragraph_spec.alignment is not None:
            paragraph.alignment = alignments[paragraph_spec.alignment]
        paragraph.paragraph_format.page_break_before = paragraph_spec.page_break_before
        paragraph.paragraph_format.keep_with_next = paragraph_spec.keep_with_next
    for table_spec in request.tables:
        all_rows = ([table_spec.headers] if table_spec.headers else []) + list(
            table_spec.rows
        )
        table = document.add_table(rows=len(all_rows), cols=len(all_rows[0]))
        if table_spec.style is not None:
            table.style = table_spec.style
        for row_index, values in enumerate(all_rows):
            for column_index, value in enumerate(values):
                cell = table.cell(row_index, column_index)
                cell.text = value
                if table_spec.headers and row_index == 0:
                    for run in cell.paragraphs[0].runs:
                        run.bold = True


def _replace_powerpoint_presentation(presentation: Any, plan: Any) -> None:
    for slide in presentation.slides:
        for shape in slide.shapes:
            _replace_powerpoint_shape(shape, plan)


def _replace_powerpoint_shape(shape: Any, plan: Any) -> None:
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            _replace_runs(paragraph.runs, plan)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    _replace_runs(paragraph.runs, plan)
    for child in getattr(shape, "shapes", ()):
        _replace_powerpoint_shape(child, plan)


def _slide_layout(presentation: Any, index: int) -> Any:
    if index >= len(presentation.slide_layouts):
        raise DocumentGenerationError(
            f"slide layout index {index} is unavailable in this template"
        )
    return presentation.slide_layouts[index]


def _append_powerpoint_content(
    presentation: Any, request: PowerPointPresentationRequest
) -> None:
    if request.title_slide is not None:
        title_spec = request.title_slide
        slide = presentation.slides.add_slide(
            _slide_layout(presentation, title_spec.layout_index)
        )
        title_shape = slide.shapes.title or _add_title_box(presentation, slide)
        title_shape.text = title_spec.title
        if title_spec.subtitle is not None:
            subtitle_shape = _first_text_shape(slide, exclude=title_shape)
            if subtitle_shape is None:
                subtitle_shape = _add_body_box(presentation, slide)
            subtitle_shape.text = title_spec.subtitle
    for slide_spec in request.slides:
        slide = presentation.slides.add_slide(
            _slide_layout(presentation, slide_spec.layout_index)
        )
        title_shape = slide.shapes.title or _add_title_box(presentation, slide)
        title_shape.text = slide_spec.title
        body_shape = _first_text_shape(slide, exclude=title_shape)
        if body_shape is None:
            body_shape = _add_body_box(presentation, slide)
        text_frame = body_shape.text_frame
        text_frame.clear()
        for index, item in enumerate(slide_spec.paragraphs):
            paragraph = (
                text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            )
            paragraph.text = item.text
            paragraph.level = item.level
            if paragraph.runs:
                run = paragraph.runs[0]
                run.font.bold = item.bold
                run.font.italic = item.italic
                if item.font_size_points is not None and _PowerPointPoints is not None:
                    run.font.size = _PowerPointPoints(item.font_size_points)


def _first_text_shape(slide: Any, *, exclude: Any) -> Any | None:
    for shape in slide.placeholders:
        if not _same_shape(shape, exclude) and getattr(shape, "has_text_frame", False):
            return shape
    for shape in slide.shapes:
        if not _same_shape(shape, exclude) and getattr(shape, "has_text_frame", False):
            return shape
    return None


def _same_shape(first: Any, second: Any) -> bool:
    return first is second or getattr(first, "_element", None) is getattr(
        second, "_element", None
    )


def _add_title_box(presentation: Any, slide: Any) -> Any:
    return slide.shapes.add_textbox(
        int(presentation.slide_width * 0.08),
        int(presentation.slide_height * 0.06),
        int(presentation.slide_width * 0.84),
        int(presentation.slide_height * 0.16),
    )


def _add_body_box(presentation: Any, slide: Any) -> Any:
    return slide.shapes.add_textbox(
        int(presentation.slide_width * 0.08),
        int(presentation.slide_height * 0.25),
        int(presentation.slide_width * 0.84),
        int(presentation.slide_height * 0.65),
    )


def _atomic_write(path: Path, content: bytes, *, overwrite: bool) -> None:
    if not overwrite:
        try:
            descriptor = os.open(
                path,
                os.O_WRONLY | os.O_CREAT | os.O_EXCL,
                0o600,
            )
        except FileExistsError as exc:
            raise ArtifactExistsError(f"artifact already exists: {path.name}") from exc
        with os.fdopen(descriptor, "wb") as destination:
            destination.write(content)
            destination.flush()
            os.fsync(destination.fileno())
        return
    descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{path.name}.", suffix=".tmp", dir=path.parent
    )
    temporary_path = Path(temporary_name)
    try:
        with os.fdopen(descriptor, "wb") as destination:
            destination.write(content)
            destination.flush()
            os.fsync(destination.fileno())
        os.replace(temporary_path, path)
    finally:
        temporary_path.unlink(missing_ok=True)
